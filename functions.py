import PyPDF2
import docx
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz  # PyMuPDF
import pdfplumber
import cv2
from PIL import Image
from io import BytesIO
import numpy as np
import pandas as pd
import easyocr  # <-- replaced pytesseract with EasyOCR

# =========================
# 🔤 EasyOCR: single shared reader (faster)
# =========================
_reader = None
def get_ocr_reader(langs=None):
    """
    Lazily initialize a single EasyOCR Reader instance.
    Example: langs=['en'] or langs=['en','ur'] for English + Urdu.
    """
    global _reader
    if _reader is None:
        _reader = easyocr.Reader(langs or ['en'])
    return _reader

# =========================
# 📄 PDF & IMAGE FUNCTIONS
# =========================
def extract_text_from_pdf(path: str) -> str:
    with open(path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

def extract_images_from_pdf(path: str) -> list:
    images = []
    doc = fitz.open(path)
    for i in range(len(doc)):
        page = doc.load_page(i)
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            images.append(Image.open(BytesIO(base_image["image"])))
    return images

def extract_tables_from_pdf(path: str) -> list:
    tables = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            for table in page.extract_tables():
                df = pd.DataFrame(table[1:], columns=table[0])
                tables.append(df)
    return tables

def preprocess_image_for_ocr(image_path: str) -> Image.Image:
    img = Image.open(image_path).convert("L")
    arr = np.array(img)
    thresh = cv2.adaptiveThreshold(
        arr, 255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        11, 2
    )
    denoised = cv2.fastNlMeansDenoising(thresh, None, 30, 7, 21)
    return Image.fromarray(denoised)

# =========================
# 🧠 OCR (EasyOCR)
# =========================
def image_ocr(image_path: str) -> str:
    """
    OCR an image from disk path using EasyOCR.
    Keeps your preprocessing; returns plain text.
    """
    processed = preprocess_image_for_ocr(image_path)
    reader = get_ocr_reader(['en'])  # add more languages as needed
    results = reader.readtext(np.array(processed), detail=0)  # list[str]
    text = "\n".join(results).strip()
    return text or "(no text found)"

def extract_text_from_image(image: Image.Image) -> str:
    """
    OCR a PIL Image in-memory using EasyOCR.
    """
    gray = image.convert("L")
    reader = get_ocr_reader(['en'])
    results = reader.readtext(np.array(gray), detail=0)
    return ("\n".join(results)).strip() or "(no text found)"

def extract_table_from_image(image: Image.Image) -> str:
    """
    OCR text from an image that may contain a table.
    NOTE: This returns unstructured text (like your previous pytesseract version).
    For structured tables, integrate a table-structure model later if needed.
    """
    try:
        reader = get_ocr_reader(['en'])
        results = reader.readtext(np.array(image), detail=0)
        return ("\n".join(results)).strip() or "(no text found)"
    except Exception as e:
        return f"Error extracting table from image: {e}"

# =========================
# 📝 DOCX FUNCTIONS
# =========================
def extract_text_from_docx(docx_path):
    try:
        document = docx.Document(docx_path)
        text = "\n".join([para.text for para in document.paragraphs])
        return text.strip()
    except Exception as e:
        return f"Error reading DOCX: {e}"

def extract_images_from_docx(docx_path):
    images = []
    try:
        with zipfile.ZipFile(docx_path, 'r') as z:
            for file in z.namelist():
                if file.startswith('word/media/'):
                    img_data = z.read(file)
                    image = Image.open(BytesIO(img_data))
                    images.append(image)
    except Exception as e:
        print(f"Error extracting images from DOCX: {e}")
    return images

def extract_tables_from_docx(docx_path):
    try:
        document = docx.Document(docx_path)
        tables = []
        for table in document.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
        return tables
    except Exception as e:
        return f"Error extracting tables from DOCX: {e}"

# =========================
# 📊 PPTX FUNCTIONS
# =========================
def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts).strip()
    except Exception as e:
        return f"Error reading PPTX: {e}"

def extract_tables_from_pptx(pptx_path):
    tables = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            # shape-based tables
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table.append(row_data)
                    tables.append(table)
            # image-based tables → OCR text (unstructured)
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_bytes = shape.image.blob
                    image = Image.open(BytesIO(img_bytes)).convert("RGB")
                    table_text = extract_table_from_image(image)
                    if table_text:
                        tables.append(table_text)
    except Exception as e:
        print(f"Error extracting tables from PPTX: {e}")
    return tables

def extract_images_from_pptx(pptx_path):
    images = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_bytes = shape.image.blob
                    image = Image.open(BytesIO(img_bytes)).convert("RGB")
                    images.append(image)
    except Exception as e:
        print(f"Error extracting images from PPTX: {e}")
    return images
